import os
import PyPDF2
import json
import traceback

import PyPDF2
import docx
from pptx import Presentation

def read_file(file):
    if file.name.endswith(".pdf"):
        try:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
            
        except Exception as e:
            raise Exception("error reading the PDF file")
        
    elif file.name.endswith(".txt"):
        return file.read().decode("utf-8")
    
    elif file.name.endswith(".docx"):
        try:
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except Exception as e:
            raise Exception("error reading the DOCX file")
    
    elif file.name.endswith(".pptx"):
        try:
            presentation = Presentation(file)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise Exception("error reading the PPTX file")
    
    else:
        raise Exception("unsupported file format; only PDF, DOCX, PPTX, and TXT files are supported")


def get_table_data(quiz_str):
    try:
        # convert the quiz from a str to dict
        print("from utils")
        print(quiz_str)
        start_index = quiz_str.find('{')
        if start_index != -1:
            quiz_str = quiz_str[start_index:]
        quiz_dict=json.loads(quiz_str)
        print(1)
        print(quiz_dict)
        quiz_table_data=[]
        
        # iterate over the quiz dictionary and extract the required information
        for key, value in quiz_dict.items():
            mcq = value["mcq"]
            options_dict = value["options"]
            correct = value["correct"]
            explanation = value["explanation"]
            
            # Prepare the row with MCQ, correct answer, and explanation
            # row = {
            #     "MCQ": mcq,
            #     "Correct": correct,
            #     "Explanation": explanation
            # }
            row = {}
            row["MCQ"] = mcq
            
            # Add each option as a separate column
            for option, option_value in options_dict.items():
                row[option] = option_value
            correct_ans = row[correct]
            row["Correct"] = correct_ans
            row["Explanation"] = explanation
            
            quiz_table_data.append(row)
        
        return quiz_table_data
        
    except Exception as e:
        traceback.print_exception(type(e), e, e.__traceback__)
        return False

